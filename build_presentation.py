"""Gera Apresentacao.pptx do projeto Guardião dos Biomas.

Uso: py build_presentation.py
Saída: Apresentacao.pptx no diretório raiz do projeto.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Paleta — mesma do PRESENTATION.md
BG_DARK = RGBColor(0x0A, 0x0E, 0x1A)
BG_PANEL = RGBColor(0x14, 0x1A, 0x2A)
TXT_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
TXT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
TXT_DIM = RGBColor(0x6B, 0x72, 0x80)
ACCENT_GREEN = RGBColor(0x66, 0xEB, 0xA0)   # Mata / ODS 15
ACCENT_BLUE = RGBColor(0x75, 0xCC, 0xFF)    # Maré / ODS 14
ACCENT_GOLD = RGBColor(0xFF, 0xC7, 0x66)    # Solis / ODS 7
ACCENT_RED = RGBColor(0xFF, 0x86, 0x6B)     # Raíz / ODS 13
ACCENT_HOT = RGBColor(0xFF, 0x4D, 0x4D)     # boss final

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ---------- helpers ----------

def add_background(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_accent_bar(slide, color=ACCENT_GREEN, top=True, height=Emu(60000)):
    if top:
        y = 0
    else:
        y = SH - height
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SW, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_text(slide, left, top, width, height, text, size=18,
             color=TXT_WHITE, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return box


def add_bullets(slide, left, top, width, height, items, size=18,
                color=TXT_WHITE, bullet_color=None, font_name="Calibri",
                line_spacing=1.35, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    bullet_color = bullet_color or ACCENT_GREEN
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(6)
        # Bullet glyph
        bullet_run = p.add_run()
        bullet_run.text = "▸  "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True
        bullet_run.font.name = font_name
        # Body
        body_run = p.add_run()
        body_run.text = item
        body_run.font.size = Pt(size)
        body_run.font.color.rgb = color
        body_run.font.name = font_name
        if bold_first and i == 0:
            body_run.font.bold = True
    return box


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_card(slide, left, top, width, height, fill=BG_PANEL,
             border_color=None, border_width=Pt(1)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border_color is not None:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        card.line.fill.background()
    card.shadow.inherit = False
    return card


def add_eyebrow(slide, top, text, color=ACCENT_GREEN):
    add_text(slide, Inches(0.7), top, Inches(12), Inches(0.4),
             text, size=14, color=color, bold=True, align=PP_ALIGN.LEFT)


def add_pagenum(slide, n, total):
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.4), Inches(0.3),
             f"{n:02d} / {total:02d}", size=10, color=TXT_DIM, align=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 14


# ---------- SLIDE 1: Capa ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GREEN, top=True)
add_accent_bar(s, ACCENT_GREEN, top=False)

# Eyebrow superior
add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.5),
         "FATEC JACAREÍ · DSM  ·  DIA DO MEIO AMBIENTE  ·  02 JUN 2026",
         size=14, color=ACCENT_GREEN, bold=True)

# Título principal
add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(1.6),
         "GUARDIÃO DOS BIOMAS", size=72, color=TXT_WHITE, bold=True,
         font_name="Calibri")

# Subtítulo
add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.7),
         "Um jogo de plataforma 2D sobre os Objetivos de Desenvolvimento Sustentável",
         size=22, color=TXT_GRAY)

# Divisor
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(5.1),
                          Inches(1.5), Emu(20000))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_GREEN
div.line.fill.background()

# Equipe pequena no rodapé
add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.4),
         "EQUIPE",
         size=12, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(1.0),
         "Alicia Silva Dias  ·  Gabrielly Neu dos Santos  ·  Manuela Lucia Lemes de Castro  ·  Pedro Claudino Nunes  ·  Gabriel Teodoro",
         size=16, color=TXT_WHITE)

add_pagenum(s, 1, TOTAL_SLIDES)
add_notes(s, "Boa tarde, somos o time do Guardião dos Biomas — um jogo de plataforma 2D "
             "que desenvolvemos pra trazer os ODS pra dentro de um game divertido. Em 10 "
             "minutos vamos mostrar o conceito, as mecânicas e como o jogo nasceu.")


# ---------- SLIDE 2: O Problema ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GOLD, top=True)
add_eyebrow(s, Inches(0.6), "01 · CONTEXTO", ACCENT_GOLD)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "Por que jogos para falar de ODS?",
         size=42, color=TXT_WHITE, bold=True)

bullets = [
    "17 ODS foram lançados pela ONU em 2015, com prazo até 2030.",
    "A maioria dos jovens conhece o nome \"ODS\", mas não conecta com o cotidiano.",
    "Educação ambiental tradicional não engaja a geração que cresceu com videogame.",
    "Jogo = ferramenta de empatia. O jogador SENTE o problema, não só lê sobre.",
]
add_bullets(s, Inches(0.7), Inches(2.5), Inches(12), Inches(4.0),
            bullets, size=22, bullet_color=ACCENT_GOLD)

# Citação no rodapé
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "\"Games são a forma mais alta de arte interativa.\"  —  Hideo Kojima",
         size=14, color=TXT_DIM)

add_pagenum(s, 2, TOTAL_SLIDES)
add_notes(s, "A ONU lançou 17 ODS em 2015. Conhecer os ODS é fácil — agir é difícil. "
             "Resolvemos que o jeito mais direto de conectar a galera da nossa idade com "
             "causa ambiental é colocando ela DENTRO do problema, controlando um personagem "
             "que vive aquilo. Daí veio o Guardião dos Biomas.")


# ---------- SLIDE 3: Conceito ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GREEN, top=True)
add_eyebrow(s, Inches(0.6), "02 · CONCEITO", ACCENT_GREEN)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "🌍  Guardião dos Biomas",
         size=42, color=TXT_WHITE, bold=True)

add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
         "Plataforma 2D pixel art (estilo retro 16-bit)",
         size=18, color=TXT_GRAY)

# 4 ODS cards
card_top = Inches(3.0)
card_h = Inches(2.4)
card_w = Inches(2.85)
card_gap = Inches(0.18)
card_start_x = Inches(0.7)

ods_data = [
    ("ODS 15", "🌳 FLORESTA", "Vida Terrestre", ACCENT_GREEN),
    ("ODS 14", "🌊 OCEANO", "Vida na Água", ACCENT_BLUE),
    ("ODS 7", "⚡ CIDADE", "Energia Limpa", ACCENT_GOLD),
    ("ODS 13", "🔥 BOSS FINAL", "Ação Climática", ACCENT_RED),
]
for i, (ods, name, sub, color) in enumerate(ods_data):
    x = card_start_x + (card_w + card_gap) * i
    add_card(s, x, card_top, card_w, card_h, border_color=color, border_width=Pt(1.5))
    add_text(s, x, card_top + Inches(0.25), card_w, Inches(0.4),
             ods, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, card_top + Inches(0.85), card_w, Inches(0.6),
             name, size=20, color=TXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, card_top + Inches(1.65), card_w, Inches(0.4),
             sub, size=14, color=TXT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.8),
         "4 fases temáticas  ·  4 personagens jogáveis  ·  inimigos personificam ameaças reais",
         size=16, color=TXT_GRAY, align=PP_ALIGN.LEFT)

add_pagenum(s, 3, TOTAL_SLIDES)
add_notes(s, "O jogo tem 4 fases. Cada uma representa um ODS específico. Em vez de só ouvir "
             "\"preserve a floresta\", o jogador encara o desmatador. Em vez de só ler "
             "\"reduza plástico no oceano\", ele recolhe lixo e separa em lixeiras corretas. "
             "Os inimigos não são monstros aleatórios — são metáforas das ameaças reais.")


# ---------- SLIDE 4: Os 4 Guardiões ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GREEN, top=True)
add_eyebrow(s, Inches(0.6), "03 · PERSONAGENS", ACCENT_GREEN)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "Escolha seu Guardião",
         size=42, color=TXT_WHITE, bold=True)

guardians = [
    ("MATA", "Mata Atlântica", "ODS 15", "Pulo +25%", ACCENT_GREEN),
    ("MARÉ", "Manguezal", "ODS 14", "Velocidade +30%", ACCENT_BLUE),
    ("SOLIS", "Cerrado", "ODS 7", "Coleta 2×", ACCENT_GOLD),
    ("RAÍZ", "Caatinga", "ODS 13", "Vida extra", ACCENT_RED),
]

card_top = Inches(2.7)
card_h = Inches(3.5)
card_w = Inches(2.85)
card_gap = Inches(0.18)
card_start_x = Inches(0.7)

for i, (name, bioma, ods, power, color) in enumerate(guardians):
    x = card_start_x + (card_w + card_gap) * i
    add_card(s, x, card_top, card_w, card_h, border_color=color, border_width=Pt(1.5))
    # Top color stripe inside card
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.15), card_top + Inches(0.15),
                                 card_w - Inches(0.3), Inches(0.07))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    add_text(s, x, card_top + Inches(0.45), card_w, Inches(0.7),
             name, size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, card_top + Inches(1.25), card_w, Inches(0.4),
             bioma, size=15, color=TXT_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, card_top + Inches(1.65), card_w, Inches(0.4),
             ods, size=13, color=TXT_DIM, align=PP_ALIGN.CENTER)
    # Power line
    add_text(s, x, card_top + Inches(2.6), card_w, Inches(0.4),
             "PODER", size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, card_top + Inches(2.95), card_w, Inches(0.4),
             power, size=16, color=TXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 4, TOTAL_SLIDES)
add_notes(s, "Cada Guardião tem uma habilidade que combina com o bioma de onde vem. "
             "Mata, da Mata Atlântica, é ágil e salta alto entre as árvores. Maré, do "
             "Manguezal, é veloz como a corrente. Solis, do Cerrado ensolarado, sente "
             "os recursos à distância. Raíz, da Caatinga resistente, aguenta um a mais "
             "— começa o jogo com vida extra. O jogador escolhe um e parte pra missão.")


# ---------- SLIDE 5: Fase 1 — Floresta ----------

def fase_slide(idx, ods_num, ods_name, title_emoji, fase_title, color,
               objetivo, inimigo, boss, fato, acao, note):
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_accent_bar(s, color, top=True)
    add_eyebrow(s, Inches(0.6), f"FASE  ·  ODS {ods_num}  ·  {ods_name.upper()}", color)
    add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.1),
             f"{title_emoji}  {fase_title}",
             size=42, color=TXT_WHITE, bold=True)

    # Cards: Objetivo, Inimigo, Boss
    card_top = Inches(2.6)
    card_h = Inches(1.55)
    card_w = Inches(3.85)
    card_gap = Inches(0.22)
    card_start_x = Inches(0.7)
    fields = [
        ("OBJETIVO", objetivo),
        ("INIMIGOS", inimigo),
        ("CHEFE", boss),
    ]
    for i, (label, body) in enumerate(fields):
        x = card_start_x + (card_w + card_gap) * i
        add_card(s, x, card_top, card_w, card_h, border_color=color, border_width=Pt(1))
        add_text(s, x + Inches(0.3), card_top + Inches(0.18), card_w - Inches(0.6),
                 Inches(0.35), label, size=12, color=color, bold=True)
        add_text(s, x + Inches(0.3), card_top + Inches(0.6),
                 card_w - Inches(0.6), Inches(0.9),
                 body, size=16, color=TXT_WHITE)

    # Fato real
    add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
             "FATO APRESENTADO AO JOGADOR", size=12, color=color, bold=True)
    add_text(s, Inches(0.7), Inches(4.95), Inches(12), Inches(1.0),
             f"\"{fato}\"", size=18, color=TXT_WHITE)

    # Ação
    add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.4),
             "AÇÃO SUGERIDA", size=12, color=color, bold=True)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
             acao, size=16, color=TXT_GRAY)

    add_pagenum(s, idx, TOTAL_SLIDES)
    add_notes(s, note)
    return s


fase_slide(5, 15, "Vida Terrestre", "🌳", "Floresta", ACCENT_GREEN,
           objetivo="Coletar 5 sementes + derrotar o Ogro Devastador.",
           inimigo="Lodos animados — criaturas da terra contaminada pelo desmatamento.",
           boss="Ogre — personificação do desmatamento brutal.",
           fato="Florestas cobrem 31% da Terra e abrigam 80% das espécies terrestres.",
           acao="Plante uma árvore nativa. Reduza consumo de papel.",
           note="A Floresta abre o jogo. O jogador coleta sementes — gesto simbólico de "
                "reflorestamento — enquanto desvia de lodos hostis. No fim, enfrenta o ogro "
                "que personifica o desmatamento. Ao terminar a fase, recebe um fato real "
                "sobre cobertura florestal e uma ação concreta que pode fazer.")


# ---------- SLIDE 6: Fase 2 — Oceano ----------

fase_slide(6, 14, "Vida na Água", "🌊", "Oceano", ACCENT_BLUE,
           objetivo="Coletar 5 itens recicláveis e depositar cada um na lixeira correta.",
           inimigo="Lesmas tóxicas — poluição visível do oceano.",
           boss="Big Zombie — vítima do plástico no mar.",
           fato="11 milhões de toneladas de plástico chegam ao oceano todo ano.",
           acao="Separe seu lixo. Pesquise a coleta da sua cidade.",
           note="A fase Oceano é a mais educativa: o jogador NÃO PODE só pegar e ignorar — "
                "tem que escolher a lixeira certa. Acertar = +1 reciclado. Errar = perde "
                "tempo. É uma mecânica que ensina separação de resíduos jogando, não "
                "palestrando. O boss aqui é literalmente uma vítima do plástico.")


# ---------- SLIDE 7: Fase 3 — Cidade ----------

fase_slide(7, 7, "Energia Limpa", "⚡", "Cidade do Amanhã", ACCENT_GOLD,
           objetivo="Coletar 5 baterias solares + derrotar o demônio industrial.",
           inimigo="Chorts vermelhos — poluição industrial personificada.",
           boss="Big Demon — última resistência do modelo poluente.",
           fato="700 milhões de pessoas vivem sem acesso à eletricidade. Energia limpa "
                "= 30% da matriz mundial e crescendo.",
           acao="Economize energia. Apoie políticas de transição energética.",
           note="A Cidade do Amanhã mostra que sustentabilidade não é antiurbana — é uma "
                "cidade que encontrou como conviver com energia limpa. Plataformas eólicas, "
                "prédios eficientes, baterias. O boss é o demônio industrial.")


# ---------- SLIDE 8: Boss Final ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_HOT, top=True)
add_accent_bar(s, ACCENT_HOT, top=False)
add_eyebrow(s, Inches(0.6), "CLIMAX  ·  ODS 13  ·  AÇÃO CLIMÁTICA", ACCENT_HOT)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.1),
         "🔥  Coração do Mundo",
         size=42, color=TXT_WHITE, bold=True)
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
         "Super Boss — \"O Vazio\"",
         size=22, color=ACCENT_HOT, bold=True)

bullets = [
    "Arena especial — silêncio dramático que dá lugar a música épica.",
    "Boss voador (cogumelo místico corrompido). 8 HP, padrão de smash.",
    "Mecânica: aguarde, esquive, contra-ataque com seed na janela vulnerável.",
    "Mensagem central: cada Guardião enfrenta o mesmo inimigo final — a indiferença.",
    "Vencer = tela de vitória com baú dourado + créditos rolando.",
]
add_bullets(s, Inches(0.7), Inches(3.0), Inches(12), Inches(3.5),
            bullets, size=20, bullet_color=ACCENT_HOT)

add_pagenum(s, 8, TOTAL_SLIDES)
add_notes(s, "O super boss é o clímax. A ODS 13 — Ação Climática — não tem inimigo único; "
             "é a soma de todos os outros. Por isso o boss final é \"O Vazio\": a "
             "indiferença coletiva. Quando o jogador vence, ele recebe a mensagem mais "
             "importante do jogo: cada ação conta, cada bioma vive. Aí vem a tela de "
             "vitória e os créditos.")


# ---------- SLIDE 9: Tecnologia ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_BLUE, top=True)
add_eyebrow(s, Inches(0.6), "04 · TECNOLOGIA", ACCENT_BLUE)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "Como construímos",
         size=42, color=TXT_WHITE, bold=True)

tech_left = [
    ("ENGINE", "Godot 4.6 — open-source, gratuita, leve."),
    ("LINGUAGEM", "GDScript — sintaxe Python-like, ideal pra prototipagem rápida."),
    ("ARQUITETURA", "State machines (player, inimigos, bosses) + signal-driven UI."),
]
tech_right = [
    ("PADRÕES", "Autoloads singletons (GameState, MusicPlayer, SfxPlayer)."),
    ("RENDERER", "Forward Mobile (Vulkan) — performance em hardware modesto."),
    ("VERSIONAMENTO", "Git + GitHub. Exportável pra Windows, Linux, Mac e Web."),
]

def tech_block(slide, x, items):
    y = Inches(2.6)
    for label, body in items:
        add_text(slide, x, y, Inches(5.8), Inches(0.4), label,
                 size=12, color=ACCENT_BLUE, bold=True)
        add_text(slide, x, y + Inches(0.45), Inches(5.8), Inches(0.8), body,
                 size=16, color=TXT_WHITE)
        y += Inches(1.45)

tech_block(s, Inches(0.7), tech_left)
tech_block(s, Inches(6.85), tech_right)

add_pagenum(s, 9, TOTAL_SLIDES)
add_notes(s, "Escolhemos a Godot 4 porque é open-source, gratuita e tem uma linguagem "
             "(GDScript) próxima do Python — então aprendemos rápido. Tudo é organizado em "
             "cenas (.tscn) que combinam visual + script. Usamos state machines pros "
             "personagens, autoloads pra estado global, e signals pro UI reagir a eventos.")


# ---------- SLIDE 10: Mecânicas ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GOLD, top=True)
add_eyebrow(s, Inches(0.6), "05 · GAMEPLAY", ACCENT_GOLD)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "Mecânicas de gameplay",
         size=42, color=TXT_WHITE, bold=True)

mecanicas = [
    ("Movimentação", "andar · agachar · deslizar · pular · pulo duplo"),
    ("Combate corpo-a-corpo", "stomp clássico — pula em cima do inimigo (caindo) o derrota"),
    ("Combate à distância", "lança sementes (X / J) — projétil que voa horizontal"),
    ("Wall slide", "desliza na parede vertical e dá wall-jump"),
    ("Combo timer", "mantenha ritmo coletando — se demorar entre coletas, perde vida"),
    ("Checkpoints", "bandeiras coloridas salvam progresso na fase"),
    ("Life pickups", "corações no meio do mapa recuperam vida"),
    ("HUD dinâmica", "vidas, ODS atual, contagem de coletáveis, barra de combo"),
]
add_bullets(s, Inches(0.7), Inches(2.5), Inches(12), Inches(4.5),
            [f"{label} — {body}" for label, body in mecanicas],
            size=18, bullet_color=ACCENT_GOLD, line_spacing=1.25)

add_pagenum(s, 10, TOTAL_SLIDES)
add_notes(s, "Não é só pular e correr. O jogador tem 8 mecânicas diferentes. Pode atacar "
             "de longe (lançando sementes) ou de cima (stomp). Cada coletável renova o "
             "combo timer. Checkpoints salvam progresso pra não punir morte sem "
             "necessidade. Corações no mapa recuperam vida.")


# ---------- SLIDE 11: Arte & Som ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_RED, top=True)
add_eyebrow(s, Inches(0.6), "06 · ARTE & SOM", ACCENT_RED)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "Identidade visual e sonora",
         size=42, color=TXT_WHITE, bold=True)

# 2 colunas
col_w = Inches(5.85)

# Arte
add_text(s, Inches(0.7), Inches(2.4), col_w, Inches(0.4),
         "🎨  ARTE", size=14, color=ACCENT_RED, bold=True)
arte = [
    "Personagens — Pixel Adventure (Pixelfrog, CC0)",
    "Inimigos & bosses — 0x72 DungeonTileset II",
    "Tilesets ambientais — Kenney Seasonal (CC0)",
    "Ícones (lixo, sementes, bateria) — Flaticon",
    "Fonte UI — Press Start 2P (OFL)",
]
add_bullets(s, Inches(0.7), Inches(2.95), col_w, Inches(3.5),
            arte, size=15, bullet_color=ACCENT_RED, line_spacing=1.25)

# Som
add_text(s, Inches(6.85), Inches(2.4), col_w, Inches(0.4),
         "🎵  SOM", size=14, color=ACCENT_RED, bold=True)
som = [
    "Menu — Cyberpunk Moonlight Sonata (Visager)",
    "Floresta — Through the Dark Forest",
    "Oceano — Cruising Down 8-bit Lane",
    "Cidade — Retro Space",
    "Boss Final — The Final Boss Battle (épico)",
    "SFX — gerados em sfxr.me (8-bit retrô)",
]
add_bullets(s, Inches(6.85), Inches(2.95), col_w, Inches(3.5),
            som, size=15, bullet_color=ACCENT_RED, line_spacing=1.25)

add_pagenum(s, 11, TOTAL_SLIDES)
add_notes(s, "Não desenhamos os personagens — usamos packs open-source de qualidade "
             "profissional. O Pixel Adventure tem 4 personagens com 11 frames de animação "
             "cada, perfeitos pros 4 Guardiões. O 0x72 deu os inimigos e bosses. A música "
             "foi curada faixa-por-faixa pra combinar com o tema de cada fase.")


# ---------- SLIDE 12: Equipe ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GREEN, top=True)
add_eyebrow(s, Inches(0.6), "07 · EQUIPE", ACCENT_GREEN)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "Quem fez",
         size=42, color=TXT_WHITE, bold=True)
add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
         "FATEC Jacareí  ·  DSM",
         size=16, color=TXT_GRAY)

team = [
    "Alicia Silva Dias",
    "Gabrielly Neu dos Santos",
    "Manuela Lucia Lemes de Castro",
    "Pedro Claudino Nunes",
    "Gabriel Teodoro",
]
# Each name as a centered card
card_top = Inches(3.2)
card_h = Inches(3.2)
card_w = Inches(2.3)
card_gap = Inches(0.15)
total_w = card_w * 5 + card_gap * 4
start_x = (SW - total_w) / 2

for i, name in enumerate(team):
    x = start_x + (card_w + card_gap) * i
    add_card(s, x, card_top, card_w, card_h, border_color=ACCENT_GREEN, border_width=Pt(1))
    # Initial in circle (placeholder for photo)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + (card_w - Inches(1.2)) / 2,
                                 card_top + Inches(0.4),
                                 Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_DARK
    circle.line.color.rgb = ACCENT_GREEN
    circle.line.width = Pt(1.5)
    initial = name[0]
    add_text(s, x, card_top + Inches(0.55), card_w, Inches(0.9),
             initial, size=44, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    # Name (wrapped)
    parts = name.split(" ")
    if len(parts) >= 2:
        line1 = parts[0]
        line2 = " ".join(parts[1:])
    else:
        line1 = name
        line2 = ""
    add_text(s, x, card_top + Inches(1.95), card_w, Inches(0.5),
             line1, size=15, color=TXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, card_top + Inches(2.4), card_w, Inches(0.5),
             line2, size=13, color=TXT_GRAY, align=PP_ALIGN.CENTER)

add_pagenum(s, 12, TOTAL_SLIDES)
add_notes(s, "Esse projeto foi feito por 5 alunos do DSM da Fatec Jacareí. Cada um teve "
             "um papel específico — programação, arte, design, testes. A força do projeto "
             "está na colaboração: ninguém fez sozinho.")


# ---------- SLIDE 13: Demo ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_HOT, top=True)
add_eyebrow(s, Inches(0.6), "08 · DEMO AO VIVO", ACCENT_HOT)
add_text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(2.0),
         "🎮  AO VIVO",
         size=120, color=TXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
         "Veja o Guardião dos Biomas em ação",
         size=22, color=TXT_GRAY, align=PP_ALIGN.CENTER)

# Roteiro discreto
add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
         "Menu  →  Escolha do Guardião  →  Floresta  →  Boss Final  →  Vitória",
         size=14, color=TXT_DIM, align=PP_ALIGN.CENTER)

add_pagenum(s, 13, TOTAL_SLIDES)
add_notes(s, "[Rode o jogo. Comente o que está fazendo enquanto joga. Roteiro sugerido:] "
             "1) Mostre o menu principal (visual + música). 2) Escolha um personagem (ex: "
             "Mata). 3) Jogue 1 minuto da Floresta — coletar, atacar inimigo, ativar "
             "checkpoint. 4) Pule pra Energy. 5) Mostre a fase final — boss voador. "
             "6) Tela de vitória + créditos. Duração total: 4-5 min.")


# ---------- SLIDE 14: Próximos passos + agradecimento ----------

s = prs.slides.add_slide(BLANK)
add_background(s)
add_accent_bar(s, ACCENT_GREEN, top=True)
add_accent_bar(s, ACCENT_GREEN, top=False)
add_eyebrow(s, Inches(0.6), "09 · FECHAMENTO", ACCENT_GREEN)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         "O que vem depois",
         size=42, color=TXT_WHITE, bold=True)

# Roadmap
add_text(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(0.4),
         "ROADMAP",
         size=13, color=ACCENT_GREEN, bold=True)
roadmap = [
    "Tradução para inglês e espanhol (alcance maior)",
    "Modo cooperativo local (2 Guardiões na mesma tela)",
    "Mais fases (ODS 6 — água, ODS 12 — consumo)",
    "Versão web (jogar no navegador, sem download)",
    "Conquistas e tabela de pontuação",
]
add_bullets(s, Inches(0.7), Inches(2.95), Inches(6.0), Inches(3.5),
            roadmap, size=15, bullet_color=ACCENT_GREEN, line_spacing=1.25)

# Repo + licença
add_text(s, Inches(7.2), Inches(2.4), Inches(5.5), Inches(0.4),
         "DISPONÍVEL EM", size=13, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.2), Inches(2.95), Inches(5.5), Inches(0.7),
         "github.com/<seu-usuário>/guardiao-dos-biomas",
         size=16, color=TXT_WHITE, bold=True)
add_text(s, Inches(7.2), Inches(3.6), Inches(5.5), Inches(0.5),
         "Licença MIT — livre pra estudar, modificar, redistribuir",
         size=12, color=TXT_GRAY)

# Mensagem final
final_card_top = Inches(5.3)
add_card(s, Inches(0.7), final_card_top, Inches(12.0), Inches(1.6),
         border_color=ACCENT_GREEN, border_width=Pt(1.5))
add_text(s, Inches(0.7), final_card_top + Inches(0.25), Inches(12.0), Inches(0.5),
         "Cada ação conta.  Cada bioma vive.",
         size=24, color=TXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), final_card_top + Inches(0.85), Inches(12.0), Inches(0.5),
         "Obrigado.",
         size=18, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_pagenum(s, 14, TOTAL_SLIDES)
add_notes(s, "A gente não terminou. Esse é só o primeiro release. Queremos adicionar mais "
             "ODS, mais fases, fazer multiplayer cooperativo. E o código está no GitHub "
             "sob MIT — quem quiser estudar, modificar, fazer fork, fica à vontade. "
             "Obrigado.")


# ---------- save ----------

out = "Apresentacao.pptx"
prs.save(out)
print(f"OK — {len(prs.slides)} slides escritos em '{out}'")
